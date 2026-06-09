#!/usr/bin/env python3
"""Build the SMA final-presentation deck.

One content spec -> two renderers:
  * build_pptx()    -> deck.pptx   (deliverable, Google-Slides compatible)
  * build_preview() -> preview/*.png (faithful PNG preview, since no libreoffice)
Both share the same inch/Pt geometry so the preview matches the pptx.
Run with the SMA env python (has pptx + Pillow + playwright fonts):
  ~/micromamba/envs/SMA/bin/python presentation/build_deck.py
"""
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "assets")
PREV = os.path.join(HERE, "preview"); os.makedirs(PREV, exist_ok=True)

# ---------- design ----------
SW, SH = 13.333, 7.5                  # slide size (inches, 16:9)
MX = 1.05                             # left/right margin (inches)
CREAM = (250, 247, 243); INK = (7, 55, 99); TEAL = (15, 175, 127)
MUTED = (120, 140, 162); LINE = (222, 214, 205); CARDBG = (255, 255, 255)
def rgb(t): return RGBColor(*t)

PXIN = 120                             # preview px per inch
TTC = "/usr/share/fonts/opentype/noto/"
_WT = {"r":"NotoSansCJK-Regular.ttc","b":"NotoSansCJK-Bold.ttc","m":"NotoSansCJK-Medium.ttc",
       "l":"NotoSansCJK-Light.ttc","dl":"NotoSansCJK-DemiLight.ttc"}
_fc = {}
def pil_font(weight, pt):
    key = (weight, pt)
    if key not in _fc:
        _fc[key] = ImageFont.truetype(TTC + _WT[weight], int(round(pt * PXIN / 72)), index=3)
    return _fc[key]
PPTX_FONT = "Noto Sans TC"
_PPTX_W = {"r":False,"b":True,"m":True,"l":False,"dl":False}   # bold flag for pptx

# ================= PIL renderer =================
def _img():
    im = Image.new("RGB", (int(SW*PXIN), int(SH*PXIN)), CREAM); return im, ImageDraw.Draw(im)
def _ptext(d, x, y, s, weight, pt, color):
    d.text((x*PXIN, y*PXIN), s, font=pil_font(weight, pt), fill=color)
def _pline(d, x, y, w, color, width=4):
    d.line((x*PXIN, y*PXIN, (x+w)*PXIN, y*PXIN), fill=color, width=width)
def _pimg(im, path, x, y, w, h):
    pic = Image.open(path).convert("RGB"); iw, ih = pic.size
    sc = min(w*PXIN/iw, h*PXIN/ih); nw, nh = int(iw*sc), int(ih*sc)
    pic = pic.resize((nw, nh)); px = int(x*PXIN + (w*PXIN-nw)/2); py = int(y*PXIN)
    ImageDraw.Draw(pic).rectangle((0,0,nw-1,nh-1), outline=LINE, width=2)
    im.paste(pic, (px, py)); return nw/PXIN, nh/PXIN

# ================= pptx renderer =================
def _set_font(run, weight, pt, color):
    run.font.size = Pt(pt); run.font.bold = _PPTX_W[weight]; run.font.color.rgb = rgb(color)
    run.font.name = PPTX_FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin","a:ea","a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", PPTX_FONT)
def _tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf
def _run(p, s, weight, pt, color):
    r = p.add_run(); r.text = s; _set_font(r, weight, pt, color); return r
def _xtext(slide, x, y, w, h, s, weight, pt, color, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = _tbox(slide, x, y, w, h, anchor); p = tf.paragraphs[0]; p.alignment = align
    _run(p, s, weight, pt, color); return tf
def _xline(slide, x, y, w, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color); sh.line.fill.background()
    sh.shadow.inherit = False
def _ximg(slide, path, x, y, w, h):
    pic = Image.open(path); iw, ih = pic.size; sc = min(w/(iw/PXIN), h/(ih/PXIN))
    nw, nh = (iw/PXIN)*sc, (ih/PXIN)*sc; px = x + (w-nw)/2
    p = slide.shapes.add_picture(path, Inches(px), Inches(y), Inches(nw), Inches(nh))
    p.line.color.rgb = rgb(LINE); p.line.width = Pt(1)
    return nw, nh

# ================= shared slide layouts =================
def kicker(R, S, text):
    R["text"](0)  # placeholder; replaced below
# We implement each slide type with a function taking a "backend" dict of primitives.

def render(slides, backend):
    t = backend
    for sp in slides:
        t["new"]()
        ty = sp["type"]
        if ty == "title":
            t["text"](MX, 1.25, 7, 0.5, sp.get("kicker",""), "m", 15, MUTED)
            yy = 1.85
            for ln in sp["title"]:
                t["text"](MX, yy, 11, 1.1, ln, "b", 46, INK); yy += 0.82
            t["line"](MX, yy+0.06, 1.25, TEAL)
            t["text"](MX, yy+0.28, 11, 0.6, sp["sub"], "dl", 19, INK)
            t["text"](MX, yy+0.78, 11, 0.5, sp["sub2"], "m", 16, MUTED)
            t["text"](MX, SH-1.0, 11, 0.5, sp["authors"], "r", 14, MUTED)
        elif ty == "content":
            _common_head(t, sp)
            bx = MX
            by = 1.55 + len(sp["headline"]) * 0.66 + 0.32   # below headline, dynamic
            if sp.get("bullets"):
                for main, sub in sp["bullets"]:
                    t["dot"](bx, by+0.10)
                    t["text"](bx+0.32, by-0.05, sp.get("bw",6.4), 0.6, main, "m", 18, INK)
                    if sub: t["text"](bx+0.32, by+0.34, sp.get("bw",6.4), 0.5, sub, "r", 13.5, MUTED)
                    by += sp.get("gap",0.92)
            if sp.get("stats"):
                sy = sp.get("stats_y", 2.7); sxn = sp.get("stats_x", 8.7)
                for num, lab in sp["stats"]:
                    t["text"](sxn, sy, 2.2, 0.7, num, "b", 34, TEAL)
                    t["text"](sxn+1.85, sy+0.12, 3.0, 0.5, lab, "m", 17, INK)
                    sy += 0.9
            if sp.get("image"):
                ix = sp.get("img_x", 7.5); iy = sp.get("img_y", 2.5)
                iw = sp.get("img_w", SW-ix-0.6); ih = sp.get("img_h", 4.4)
                t["img"](os.path.join(ASSET, sp["image"]), ix, iy, iw, ih)
            if sp.get("foot"):
                t["text"](MX, SH-0.85, 11.2, 0.5, sp["foot"], "r", 12.5, MUTED)
        elif ty == "bignum":
            t["text"](MX, 1.3, 8, 0.5, sp.get("kicker",""), "m", 15, TEAL)
            t["text"](MX, 2.4, 12, 2.2, sp["num"], "b", 120, INK)
            t["line"](MX, 4.55, 1.25, TEAL)
            t["text"](MX, 4.8, 10.5, 1.2, sp["caption"], "m", 22, INK)
            if sp.get("foot"): t["text"](MX, SH-0.9, 11, 0.5, sp["foot"], "r", 13, MUTED)
        elif ty == "columns":
            _common_head(t, sp)
            n = len(sp["cols"]); gw = (SW-2*MX-0.6*(n-1))/n
            for i,(ctitle,items) in enumerate(sp["cols"]):
                cx = MX + i*(gw+0.6)
                t["line"](cx, 2.75, 0.7, TEAL)
                t["text"](cx, 2.95, gw, 0.6, ctitle, "b", 21, INK)
                yy = 3.7
                for it in items:
                    t["dot"](cx, yy+0.08, small=True)
                    t["text"](cx+0.28, yy-0.03, gw-0.3, 0.8, it, "r", 15.5, INK); yy += 0.78
        elif ty == "flow":
            _common_head(t, sp)
            stages = sp["stages"]; n=len(stages)
            bw=2.55; gap=(SW-2*MX-bw*n)/(n-1); y=3.3; bh=1.5
            cx=MX
            centers=[]
            for i,(lab,subs) in enumerate(stages):
                t["box"](cx, y, bw, bh, TEAL if i in (1,) else INK, lab, subs)
                centers.append((cx,cx+bw))
                cx+=bw+gap
            for i in range(n-1):
                t["arrow"](centers[i][1], centers[i+1][0], y+bh/2)
            if sp.get("foot"): t["text"](MX, SH-0.95, 11.5, 0.6, sp["foot"], "r", 13, MUTED)

def _common_head(t, sp):
    t["text"](MX, 1.05, 9, 0.5, sp.get("kicker",""), "m", 15, TEAL)
    yy = 1.55
    for ln in sp["headline"]:
        t["text"](MX, yy, 11.4, 0.8, ln, "b", 33, INK); yy += 0.66

# ---- PIL backend ----
def pil_backend():
    state = {}
    def new():
        im, d = _img(); state["im"]=im; state["d"]=d; state["list"].append(im)
    def text(x,y,w,h,s,wt,pt,c): _ptext(state["d"], x, y, s, wt, pt, c)
    def line(x,y,w,c): _pline(state["d"], x, y, w, c)
    def img(p,x,y,w,h): return _pimg(state["im"], p, x, y, w, h)
    def dot(x,y,small=False):
        r=(5 if small else 7); cx,cy=x*PXIN,y*PXIN
        state["d"].ellipse((cx,cy,cx+r*2,cy+r*2), fill=TEAL)
    def box(x,y,w,h,col,lab,subs):
        d=state["d"]; d.rounded_rectangle((x*PXIN,y*PXIN,(x+w)*PXIN,(y+h)*PXIN), radius=8,
            outline=col, width=3, fill=(col[0],col[1],col[2]) if False else CARDBG)
        d.text((x*PXIN+14, y*PXIN+14), lab, font=pil_font("b",17), fill=col)
        yy=y*PXIN+46
        for s in subs:
            d.text((x*PXIN+14, yy), s, font=pil_font("r",12), fill=MUTED); yy+=22
    def arrow(x1,x2,y):
        d=state["d"]; d.line((x1*PXIN+4,y*PXIN,x2*PXIN-4,y*PXIN), fill=INK, width=3)
        ax=x2*PXIN-4; d.polygon([(ax-8,y*PXIN-6),(ax,y*PXIN),(ax-8,y*PXIN+6)], fill=INK)
    state["list"]=[]
    return dict(new=new,text=text,line=line,img=img,dot=dot,box=box,arrow=arrow,_state=state)

# ---- pptx backend ----
def pptx_backend(prs):
    state={"slide":None}
    blank=prs.slide_layouts[6]
    def new():
        s=prs.slides.add_slide(blank); s.background.fill.solid()
        s.background.fill.fore_color.rgb=rgb(CREAM); state["slide"]=s
    def text(x,y,w,h,s,wt,pt,c): _xtext(state["slide"],x,y,w,h,s,wt,pt,c)
    def line(x,y,w,c): _xline(state["slide"],x,y,w,c)
    def img(p,x,y,w,h): return _ximg(state["slide"],p,x,y,w,h)
    def dot(x,y,small=False):
        d=0.1 if small else 0.13
        sh=state["slide"].shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(TEAL); sh.line.fill.background(); sh.shadow.inherit=False
    def box(x,y,w,h,col,lab,subs):
        sh=state["slide"].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(CARDBG); sh.line.color.rgb=rgb(col); sh.line.width=Pt(1.5)
        sh.shadow.inherit=False
        tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.12); tf.margin_top=Inches(0.1)
        p=tf.paragraphs[0]; _run(p,lab,"b",15,col)
        for s in subs:
            pp=tf.add_paragraph(); _run(pp,s,"r",10.5,MUTED)
    def arrow(x1,x2,y):
        sh=state["slide"].shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1+0.03),Inches(y-0.08),Inches(x2-x1-0.06),Inches(0.16))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(INK); sh.line.fill.background(); sh.shadow.inherit=False
    return dict(new=new,text=text,line=line,img=img,dot=dot,box=box,arrow=arrow)

# ================= content =================
from deck_content import SLIDES

def main():
    # preview
    pb=pil_backend(); render(SLIDES, pb)
    for i,im in enumerate(pb["_state"]["list"],1):
        im.save(os.path.join(PREV, f"s{i:02d}.png"))
    # pptx
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    render(SLIDES, pptx_backend(prs))
    prs.save(os.path.join(HERE,"deck.pptx"))
    print(f"OK  {len(SLIDES)} slides -> deck.pptx + preview/s01..s{len(SLIDES):02d}.png")

if __name__=="__main__":
    main()
